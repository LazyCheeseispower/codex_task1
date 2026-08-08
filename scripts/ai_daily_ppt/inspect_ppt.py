#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Inspect a PPTX layout: shape geometry, text metrics and overflow estimates.

This is the geometry half of "actually understanding the deck": it dumps every
shape that matters, estimates how many lines its text needs, and flags text
boxes whose estimated content cannot fit their frame. It also reports pairwise
overlaps between text-bearing shapes so layout bugs are caught before render.
"""

import argparse
import json
import math
import sys

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu


EMU_PER_INCH = 914400


def shape_type(shape):
    if shape.shape_type is None:
        return "unknown"
    return str(shape.shape_type).split("(")[0].strip()


def est_text_width(ch, size_pt):
    em = size_pt / 72.0
    if ord(ch) > 0x2E80 or ch in "，。：；！？、（）《》「」『』·":
        return em
    if ch == " ":
        return em * 0.30
    if ch.isascii():
        return em * 0.55
    return em * 0.72


def est_lines(text, size_pt, width_in):
    if not text:
        return 0
    usable = max(width_in - 0.08, 0.2)
    line_width = 0.0
    lines = 1
    for ch in text:
        w = est_text_width(ch, size_pt)
        if line_width + w > usable and line_width > 0:
            lines += 1
            line_width = w
        else:
            line_width += w
    return lines


def text_frame_info(shape):
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame
    paras = []
    max_size = 0
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            size = r.font.size.pt if r.font.size is not None else 0
            max_size = max(max_size, size)
            color = None
            try:
                if r.font.color and r.font.color.type is not None:
                    color = str(r.font.color.rgb)
            except Exception:
                color = None
            runs.append({
                "text": r.text,
                "size": size,
                "bold": bool(r.font.bold),
                "color": color,
            })
        paras.append({"text": p.text, "runs": runs})
    auto = tf.auto_size
    auto_name = "NONE"
    if auto is MSO_AUTO_SIZE.NONE:
        auto_name = "NONE"
    elif auto is MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
        auto_name = "SHAPE_TO_FIT_TEXT"
    elif auto is MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
        auto_name = "TEXT_TO_FIT_SHAPE"
    text = tf.text.strip()
    if not text:
        return {"text": "", "paragraphs": paras, "max_size": 0}
    size = max_size or 12
    width_in = shape.width / EMU_PER_INCH if shape.width else 0
    height_in = shape.height / EMU_PER_INCH if shape.height else 0
    need_lines = sum(
        max(est_lines(p["text"], max((r["size"] for r in p["runs"]), default=size) or size, width_in), 1)
        for p in paras
    )
    line_height_in = size * 1.25 / 72.0
    need_in = need_lines * line_height_in
    return {
        "text": text,
        "paragraphs": paras,
        "max_size": size,
        "word_wrap": bool(tf.word_wrap),
        "auto_size": auto_name,
        "anchor": str(tf.vertical_anchor).split("(")[0].strip() if tf.vertical_anchor is not None else None,
        "est_lines": need_lines,
        "est_height_in": round(need_in, 3),
        "frame_height_in": round(height_in, 3),
        "overflow": round(need_in - height_in, 3) if height_in else None,
    }


def rect(shape):
    left = shape.left if shape.left is not None else 0
    top = shape.top if shape.top is not None else 0
    width = shape.width if shape.width is not None else 0
    height = shape.height if shape.height is not None else 0
    return left, top, width, height


def rect_in(shape):
    left, top, width, height = rect(shape)
    return {
        "left": round(left / EMU_PER_INCH, 3),
        "top": round(top / EMU_PER_INCH, 3),
        "width": round(width / EMU_PER_INCH, 3),
        "height": round(height / EMU_PER_INCH, 3),
    }


def overlaps(a, b):
    ax0, ay0, aw, ah = rect(a)
    bx0, by0, bw, bh = rect(b)
    ax1, ay1 = ax0 + aw, ay0 + ah
    bx1, by1 = bx0 + bw, by0 + bh
    ox = min(ax1, bx1) - max(ax0, bx0)
    oy = min(ay1, by1) - max(ay0, by0)
    if ox <= 0 or oy <= 0:
        return None
    return {
        "overlap_in": round((ox / EMU_PER_INCH) * (oy / EMU_PER_INCH), 4),
        "dx_in": round(ox / EMU_PER_INCH, 3),
        "dy_in": round(oy / EMU_PER_INCH, 3),
    }


def dump_slide(index, slide, slide_w, slide_h, issues):
    shapes = []
    for shape in slide.shapes:
        info = {
            "name": shape.name,
            "type": shape_type(shape),
            "rect": rect_in(shape),
        }
        if shape.has_text_frame:
            info["text_frame"] = text_frame_info(shape)
        if getattr(shape, "has_table", False) and shape.has_table:
            info["table"] = {
                "rows": len(shape.table.rows),
                "cols": len(shape.table.columns),
            }
        shapes.append(info)

        # Bounds check.
        left, top, width, height = rect(shape)
        tol = Emu(int(0.02 * EMU_PER_INCH))
        if left < -tol or top < -tol or left + width > slide_w + tol or top + height > slide_h + tol:
            issues.append("slide %d: %r out of bounds" % (index + 1, shape.name))

        tf = info.get("text_frame")
        if tf and tf.get("text"):
            overflow = tf.get("overflow")
            if overflow is not None and overflow > 0.02:
                issues.append(
                    "slide %d: %r estimated %.1f lines / %.2fin needs %.2fin"
                    % (index + 1, shape.name, tf["est_lines"], tf["frame_height_in"], tf["est_height_in"])
                )

    # Overlap text-bearing shapes with non-trivial content.
    text_shapes = [
        sh for sh in slide.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    ]
    for i in range(len(text_shapes)):
        for j in range(i + 1, len(text_shapes)):
            a, b = text_shapes[i], text_shapes[j]
            ov = overlaps(a, b)
            if ov:
                issues.append(
                    "slide %d: text shapes %r and %r overlap %.3fin x %.3fin"
                    % (index + 1, a.name, b.name, ov["dx_in"], ov["dy_in"])
                )

    return shapes


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--report", required=True)
    ap.add_argument("--slides", default="")
    args = ap.parse_args()

    prs = Presentation(args.pptx)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    wanted = set()
    if args.slides:
        wanted = {int(x) - 1 for x in args.slides.split(",") if x.strip()}
    issues = []
    out = {
        "slide_width_in": round(slide_w / EMU_PER_INCH, 3),
        "slide_height_in": round(slide_h / EMU_PER_INCH, 3),
        "slides": [],
        "issues": issues,
    }
    for i, slide in enumerate(prs.slides):
        if wanted and i not in wanted:
            continue
        out["slides"].append({
            "index": i,
            "layout": slide.slide_layout.name,
            "shape_count": len(slide.shapes),
            "shapes": dump_slide(i, slide, slide_w, slide_h, issues),
        })
    with open(args.report, "w", encoding="utf-8") as f:
        json.dump(out, f, ensure_ascii=False, indent=2)
    print("slides dumped: %d, issues: %d" % (len(out["slides"]), len(issues)))
    for issue in issues[:80]:
        print(" -", issue)


if __name__ == "__main__":
    main()
