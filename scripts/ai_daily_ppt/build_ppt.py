#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build a brand-consistent AI daily brief from the Zhende template.

The generator never redraws the template. It reuses the template's own cover,
global TOC, chapter divider/TOC, content skeleton and back-cover slides, then
overlays an adaptive card system so every run keeps the same visual language:
4 items use a 2x2 grid, 3 items use three balanced columns, 2 items use two
larger cards and 1 item uses a banner card. Every text frame enables PowerPoint
auto-shrink so long headlines cannot spill out of their card.
"""

import argparse
import copy
import json
import sys
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


BRAND_BLUE = RGBColor(0x15, 0x3C, 0x86)
ACCENT_BLUE = RGBColor(0x29, 0x3F, 0x8E)
DARK_TEXT = RGBColor(0x22, 0x2B, 0x36)
BODY_TEXT = RGBColor(0x3A, 0x44, 0x52)
META_TEXT = RGBColor(0x8A, 0x94, 0xA3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_LINE = RGBColor(0xD9, 0xE1, 0xEE)
LIGHT_LINE = RGBColor(0xE4, 0xEA, 0xF3)
FONT = "微软雅黑"

MAX_ITEMS_PER_CHAPTER = 4
MAX_CHAPTERS = 6

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.76
GAP = 0.30
GRID_TOP = 1.55
GRID_BOTTOM = 6.75


TIERS = {
    "compact": {
        "title": 14.5, "summary": 10.5, "meta": 9.0, "chip": 10.5,
        "chip_h": 0.34, "chip_top": 0.22, "title_top": 0.66, "title_h": 0.58,
        "summary_top": 1.32, "summary_h": 0.40, "meta_bottom": 0.44,
    },
    "medium": {
        "title": 14.0, "summary": 10.5, "meta": 9.0, "chip": 10.5,
        "chip_h": 0.34, "chip_top": 0.24, "title_top": 0.68, "title_h": 0.68,
        "summary_top": 1.46, "summary_h": 1.09, "meta_bottom": 0.44,
    },
    "large": {
        "title": 16.0, "summary": 11.5, "meta": 9.5, "chip": 11.0,
        "chip_h": 0.38, "chip_top": 0.28, "title_top": 0.76, "title_h": 0.76,
        "summary_top": 1.64, "summary_h": 1.47, "meta_bottom": 0.48,
    },
    "banner": {
        "title": 18.0, "summary": 12.0, "meta": 10.0, "chip": 11.0,
        "chip_h": 0.40, "chip_top": 0.32, "title_top": 0.88, "title_h": 0.82,
        "summary_top": 1.88, "summary_h": 0.61, "meta_bottom": 0.50,
    },
}


def _apply_typeface(rPr, name):
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            anchor = rPr.find(qn("a:sym"))
            if anchor is None:
                anchor = rPr.find(qn("a:hlinkClick"))
            if anchor is None:
                anchor = rPr.find(qn("a:hlinkMouseOver"))
            if anchor is None:
                anchor = rPr.find(qn("a:rtl"))
            if anchor is not None:
                anchor.addprevious(el)
            else:
                rPr.append(el)
        el.set("typeface", name)


def set_run_font(run, size=None, bold=None, color=None, name=None):
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = color
    if name:
        f.name = name
        _apply_typeface(run._r.get_or_add_rPr(), name)


def set_existing_text(shape, text, size=None, bold=None, color=None, name=None,
                      align=None, anchor=None, line_spacing=1.08):
    """Replace a template text box with a single clean run and auto-shrink."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    p = tf.paragraphs[0]
    for run in list(p.runs):
        run._r.getparent().remove(run._r)
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color, name=name)
    if align is not None:
        p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(0)
    return run


def set_cell_text(cell, text):
    tf = cell.text_frame
    p = tf.paragraphs[0]
    runs = list(p.runs)
    if runs:
        run = runs[0]
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
    run.text = text
    if not runs:
        set_run_font(run, size=10.5, color=DARK_TEXT, name=FONT)
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)


def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def find_slide(prs, predicate, label):
    for slide in prs.slides:
        if predicate(slide):
            return slide
    raise ValueError("template slide not found: %s" % label)


def duplicate_slide(prs, source):
    new_slide = prs.slides.add_slide(source.slide_layout)
    sp_tree = new_slide.shapes._spTree
    for shape in source.shapes:
        sp_tree.append(copy.deepcopy(shape._element))
    _strip_custom_data(sp_tree)
    _remap_image_rels(source, new_slide)
    return new_slide


def _strip_custom_data(sp_tree):
    for el in list(sp_tree.iter(qn("p:custDataLst"))):
        el.getparent().remove(el)


def _remap_image_rels(source, new_slide):
    for blip in new_slide.shapes._spTree.iter(qn("a:blip")):
        embed = blip.get(qn("r:embed"))
        if embed:
            try:
                part = source.part.related_part(embed)
            except KeyError:
                continue
            image_part, new_rId = new_slide.part.get_or_add_image_part(BytesIO(part.blob))
            blip.set(qn("r:embed"), new_rId)
            continue
        link = blip.get(qn("r:link"))
        if link:
            try:
                part = source.part.related_part(link)
            except KeyError:
                continue
            new_rId = new_slide.part.relate_to(part, part.reltype)
            blip.set(qn("r:link"), new_rId)


def drop_original_slides(prs, rids):
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        if rId in rids:
            prs.part.drop_rel(rId)
            sldIdLst.remove(sldId)


def add_textbox(slide, name, left, top, width, height, text, size,
                bold=False, color=None, align=None, anchor=None,
                line_spacing=1.08, auto_size=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    if auto_size:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align if align is not None else PP_ALIGN.LEFT
    p.line_spacing = line_spacing
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    set_run_font(run, size=size, bold=bold, color=color, name=FONT)
    return box


def add_rect(slide, name, left, top, width, height, color, line=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = CARD_LINE
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rounded_rect(slide, name, left, top, width, height, fill=WHITE,
                     line_color=CARD_LINE, radius=0.08):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    return shape


def add_card(slide, prefix, item_no, item, left, top, width, height, tier):
    cfg = TIERS[tier]
    card = add_rounded_rect(slide, prefix, left, top, width, height,
                            fill=WHITE, line_color=CARD_LINE, radius=0.055)
    add_rect(slide, prefix + "_accent", left, top + 0.18, 0.08, height - 0.36, BRAND_BLUE)

    chip = add_rounded_rect(
        slide, prefix + "_chip", left + 0.30, top + cfg["chip_top"],
        1.30, cfg["chip_h"], fill=BRAND_BLUE, line_color=None, radius=0.5,
    )
    add_textbox(
        slide, prefix + "_chip_text", left + 0.30, top + cfg["chip_top"],
        1.30, cfg["chip_h"], "要点 %02d" % item_no, cfg["chip"], bold=True,
        color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0, auto_size=False,
    )

    add_textbox(
        slide, prefix + "_title", left + 0.30, top + cfg["title_top"],
        width - 0.60, cfg["title_h"], item["title"], cfg["title"], bold=True,
        color=BRAND_BLUE, anchor=MSO_ANCHOR.TOP, line_spacing=1.05,
    )
    add_textbox(
        slide, prefix + "_summary", left + 0.30, top + cfg["summary_top"],
        width - 0.60, cfg["summary_h"], item["summary"], cfg["summary"],
        color=BODY_TEXT, anchor=MSO_ANCHOR.TOP, line_spacing=1.12,
    )
    meta_text = "%s · %s" % (item.get("source", ""), item.get("meta", ""))
    add_textbox(
        slide, prefix + "_meta", left + 0.30, top + height - cfg["meta_bottom"] - 0.28,
        width - 0.60, 0.28, meta_text, cfg["meta"], color=META_TEXT,
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0,
    )
    return card


def add_content_header(slide, date_text, chapter_title, page_no, page_total):
    kicker = "AI 动态日报 · %s · 第 %d 页 / 共 %d 页" % (date_text, page_no, page_total)
    add_textbox(slide, "page_kicker", MARGIN, 0.38, 10.6, 0.30, kicker, 10,
                color=META_TEXT, line_spacing=1.0)
    add_textbox(slide, "page_title", MARGIN, 0.70, 9.0, 0.56, chapter_title, 26,
                bold=True, color=BRAND_BLUE, line_spacing=1.0)
    add_rect(slide, "accent", MARGIN + 0.02, 1.34, 1.10, 0.07, ACCENT_BLUE)
    add_rect(slide, "page_rule", MARGIN, 1.52, SLIDE_W - 2 * MARGIN, 0.012, LIGHT_LINE)


def add_content_page(prs, content_source, date_text, chapter_no, chapter, page_no,
                     page_total, chunk, global_card_start):
    slide = duplicate_slide(prs, content_source)
    add_content_header(slide, date_text, chapter["title"], page_no, page_total)

    n = len(chunk)
    area_h = GRID_BOTTOM - GRID_TOP
    if n == 4:
        tier = "compact"
        card_w = (SLIDE_W - 2 * MARGIN - GAP) / 2
        card_h = (area_h - GAP) / 2
        top = GRID_TOP
        positions = [
            (MARGIN + col * (card_w + GAP), top + row * (card_h + GAP))
            for row in range(2) for col in range(2)
        ]
    elif n == 3:
        tier = "medium"
        card_w = (SLIDE_W - 2 * MARGIN - 2 * GAP) / 3
        card_h = 3.35
        top = GRID_TOP + (area_h - card_h) / 2
        positions = [(MARGIN + col * (card_w + GAP), top) for col in range(3)]
    elif n == 2:
        tier = "large"
        card_w = (SLIDE_W - 2 * MARGIN - GAP) / 2
        card_h = 3.95
        top = GRID_TOP + (area_h - card_h) / 2
        positions = [(MARGIN + col * (card_w + GAP), top) for col in range(2)]
    else:
        tier = "banner"
        card_w = SLIDE_W - 2 * MARGIN
        card_h = 3.35
        top = GRID_TOP + (area_h - card_h) / 2
        positions = [(MARGIN, top)]

    for idx, (left, top) in enumerate(positions):
        prefix = "card_%04d" % (global_card_start + idx)
        add_card(slide, prefix, idx + 1, chunk[idx], left, top, card_w, card_h, tier)
    return slide


def add_agenda_cards(toc, chapter):
    items = chapter["items"][:MAX_ITEMS_PER_CHAPTER]
    n = len(items)
    left0 = 0.95
    total_w = 12.40 - left0
    positions = []
    if n == 4:
        gap = 0.24
        w = (total_w - gap) / 2
        h = 0.98
        tops = [4.48, 5.56]
        xs = [left0, left0 + w + gap]
        positions = [(x, y, w, h) for y in tops for x in xs]
    elif n == 3:
        gap = 0.14
        w = total_w
        h = 0.72
        positions = [(left0, 4.42 + i * (h + gap), w, h) for i in range(3)]
    elif n == 2:
        gap = 0.16
        w = total_w
        h = 0.92
        positions = [(left0, 4.55 + i * (h + gap), w, h) for i in range(2)]
    else:
        w = total_w
        h = 1.20
        positions = [(left0, 4.80, w, h)]

    for idx, (left, top, w, h) in enumerate(positions):
        item = items[idx]
        add_rounded_rect(toc, "agenda_%02d" % (idx + 1), left, top, w, h,
                         fill=WHITE, line_color=CARD_LINE, radius=0.10)
        badge_h = min(0.46, h - 0.14)
        badge = add_rounded_rect(
            toc, "agenda_%02d_badge" % (idx + 1), left + 0.18, top + (h - badge_h) / 2,
            badge_h, badge_h, fill=BRAND_BLUE, line_color=None, radius=0.16,
        )
        add_textbox(
            toc, "agenda_%02d_num" % (idx + 1), left + 0.18, top + (h - badge_h) / 2,
            badge_h, badge_h, "%02d" % (idx + 1), 13, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0,
            auto_size=False,
        )
        add_textbox(
            toc, "agenda_%02d_title" % (idx + 1), left + 0.84, top + 0.06,
            w - 1.02, h - 0.12, item["title"], 16, color=DARK_TEXT,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05,
        )


def set_cover(cover, meta):
    publisher = find_shape(cover, "文本框 13")
    publisher.height = Inches(0.42)
    set_existing_text(publisher, meta["publisher"],
                      size=24, bold=False, name="思源黑体 CN Light")
    title = find_shape(cover, "文本框 12")
    title.top = Inches(4.32)
    set_existing_text(title, meta["title"],
                      size=38, bold=False, color=BRAND_BLUE, name="思源黑体 CN Medium")
    set_existing_text(find_shape(cover, "文本框 1"), meta["date"],
                      size=16, bold=False, name="思源黑体 CN Regular")
    table = find_shape(cover, "表格 7").table
    values = {
        (0, 1): meta["classification"],
        (0, 3): meta["confidentiality"],
        (1, 1): meta["department"],
        (1, 3): meta["author"],
        (2, 1): meta["created"],
        (2, 3): meta["scope"],
    }
    for (row, col), value in values.items():
        set_cell_text(table.cell(row, col), value)


def set_toc_entry(box, idx, title):
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    p = tf.paragraphs[0]
    for run in list(p.runs):
        run._r.getparent().remove(run._r)
    r1 = p.add_run()
    r1.text = "%d.  " % idx
    set_run_font(r1, size=20, bold=True, color=ACCENT_BLUE, name=FONT)
    r2 = p.add_run()
    r2.text = title
    set_run_font(r2, size=20, bold=False, color=DARK_TEXT, name=FONT)
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.0
    p.space_after = Pt(0)


def set_global_toc(toc, chapters):
    boxes = ["TextBox 2", "TextBox 3", "TextBox 4", "TextBox 6", "TextBox 7", "TextBox 8"]
    for idx, chapter in enumerate(chapters):
        box = find_shape(toc, boxes[idx])
        box.left = Inches(4.72)
        box.top = Inches(1.94 + idx * 0.78)
        box.width = Inches(6.4)
        box.height = Inches(0.62)
        set_toc_entry(box, idx + 1, chapter["title"])
    for name in boxes[len(chapters):]:
        sh = find_shape(toc, name)
        if sh is not None:
            sh._element.getparent().remove(sh._element)


def set_divider(divider, chapter_no, chapter):
    number_box = find_shape(divider, "TextBox 1")
    if number_box is None:
        number_box = find_shape(divider, "TextBox 5")
    set_existing_text(number_box, "%02d" % chapter_no, size=64, bold=True,
                      color=ACCENT_BLUE, name=FONT, anchor=MSO_ANCHOR.MIDDLE)
    title_box = find_shape(divider, "TextBox 4")
    title_box.left = Inches(8.35)
    title_box.top = Inches(3.95)
    title_box.width = Inches(4.85)
    title_box.height = Inches(0.42)
    set_existing_text(title_box, chapter["title"], size=20, bold=False,
                      color=DARK_TEXT, name=FONT, anchor=MSO_ANCHOR.MIDDLE)
    statement_box = find_shape(divider, "TextBox 6")
    statement_box.left = Inches(8.35)
    statement_box.top = Inches(4.48)
    statement_box.width = Inches(4.85)
    statement_box.height = Inches(1.40)
    set_existing_text(statement_box, chapter["statement"], size=20, bold=False,
                      color=BODY_TEXT, name=FONT, anchor=MSO_ANCHOR.TOP,
                      line_spacing=1.15)


def set_chapter_toc(toc, chapter_no, chapter):
    number_box = find_shape(toc, "TextBox 17")
    number_box.top = Inches(2.78)
    number_box.height = Inches(0.50)
    set_existing_text(number_box, "%02d" % chapter_no, size=28, bold=True,
                      color=ACCENT_BLUE, name=FONT, anchor=MSO_ANCHOR.MIDDLE)
    title_box = find_shape(toc, "TextBox 15")
    title_box.left = Inches(0.74)
    title_box.top = Inches(3.30)
    title_box.width = Inches(8.2)
    title_box.height = Inches(0.66)
    set_existing_text(title_box, chapter["title"], size=32, bold=False,
                      color=DARK_TEXT, name=FONT, anchor=MSO_ANCHOR.MIDDLE)
    for name in ("TextBox 18", "TextBox 43", "TextBox 44", "TextBox 45"):
        sh = find_shape(toc, name)
        if sh is not None:
            sh._element.getparent().remove(sh._element)
    add_agenda_cards(toc, chapter)


def build(prs, content):
    meta = content["meta"]
    chapters = content["chapters"]
    if len(chapters) > MAX_CHAPTERS:
        raise ValueError("最多支持 %d 章" % MAX_CHAPTERS)
    for chapter in chapters:
        if not 1 <= len(chapter["items"]) <= MAX_ITEMS_PER_CHAPTER:
            raise ValueError("每章 1-%d 条：%s" % (MAX_ITEMS_PER_CHAPTER, chapter["title"]))

    cover_src = find_slide(prs, lambda s: find_shape(s, "表格 7") is not None, "cover")
    global_toc_src = find_slide(prs, lambda s: find_shape(s, "TextBox 2") is not None, "global toc")
    divider_src = find_slide(
        prs,
        lambda s: s.slide_layout.name == "2_Custom Layout" and find_shape(s, "TextBox 1") is not None
        and find_shape(s, "TextBox 4") is not None,
        "divider",
    )
    toc_src = find_slide(
        prs,
        lambda s: s.slide_layout.name == "2_Custom Layout" and find_shape(s, "TextBox 17") is not None
        and find_shape(s, "TextBox 15") is not None,
        "chapter toc",
    )
    content_src = find_slide(
        prs,
        lambda s: s.slide_layout.name == "Picture with Caption" and len(s.shapes) == 0,
        "content skeleton",
    )
    back_src = find_slide(prs, lambda s: find_shape(s, "文本框 3") is not None, "back cover")

    original_rids = [sldId.get(qn("r:id")) for sldId in prs.slides._sldIdLst]

    cover = duplicate_slide(prs, cover_src)
    set_cover(cover, meta)
    global_toc = duplicate_slide(prs, global_toc_src)
    set_global_toc(global_toc, chapters)

    global_card = 0
    for chapter_no, chapter in enumerate(chapters, start=1):
        divider = duplicate_slide(prs, divider_src)
        set_divider(divider, chapter_no, chapter)
        toc = duplicate_slide(prs, toc_src)
        set_chapter_toc(toc, chapter_no, chapter)
        chunks = [
            chapter["items"][i:i + MAX_ITEMS_PER_CHAPTER]
            for i in range(0, len(chapter["items"]), MAX_ITEMS_PER_CHAPTER)
        ]
        for page_no, chunk in enumerate(chunks, start=1):
            add_content_page(
                prs, content_src, meta["date"], chapter_no, chapter,
                page_no, len(chunks), chunk, global_card,
            )
            global_card += len(chunk)

    duplicate_slide(prs, back_src)
    drop_original_slides(prs, original_rids)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True)
    ap.add_argument("--content", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()
    with open(args.content, encoding="utf-8") as f:
        content = json.load(f)
    prs = Presentation(args.template)
    build(prs, content)
    prs.save(args.out)
    print("saved %s (%d slides)" % (args.out, len(prs.slides._sldIdLst)))


if __name__ == "__main__":
    main()
