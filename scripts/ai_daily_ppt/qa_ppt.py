#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Structural QA for the v4 AI daily brief.

The deck is fixed at 15 slides: cover, global TOC, 4 chapters x (chapter page
+ 2 detail pages), back cover.  The chapter page carries a large brand image
and two preview cards; each detail page carries one picture plus a left text
panel.  This script writes a machine-readable report used by pixel QA.
"""

import argparse
import json
import sys

from pptx import Presentation
from pptx.util import Inches


CHAPTER_COUNT = 4
ITEMS_PER_CHAPTER = 2
EXPECTED_SLIDES = 15
CHAPTER_TITLES = ["模型与发布", "产品与工具", "行业与公司", "政策与观点"]

MAX_TITLE_CHARS = 30
MAX_SUMMARY_CHARS = 120
MAX_KEY_POINT_CHARS = 40
MAX_IMPACT_CHARS = 100
EMU_PER_INCH = 914400


def shape_text(shape):
    if shape is None or not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def find(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def rect_inches(shape):
    return {
        "left": round((shape.left or 0) / EMU_PER_INCH, 3),
        "top": round((shape.top or 0) / EMU_PER_INCH, 3),
        "width": round((shape.width or 0) / EMU_PER_INCH, 3),
        "height": round((shape.height or 0) / EMU_PER_INCH, 3),
    }


def picture_count(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return sum(1 for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)


def check_bounds(shape, slide_w, slide_h, issues, where):
    left = shape.left if shape.left is not None else 0
    top = shape.top if shape.top is not None else 0
    width = shape.width if shape.width is not None else 0
    height = shape.height if shape.height is not None else 0
    tol = Inches(0.02)
    if left < -tol or top < -tol or left + width > slide_w + tol or top + height > slide_h + tol:
        issues.append("%s shape %r out of bounds" % (where, shape.name))


def validate_content(content, issues):
    meta = content.get("meta", {})
    for key in ("title", "publisher", "date", "classification", "confidentiality",
                "department", "author", "created", "scope"):
        if not str(meta.get(key, "")).strip():
            issues.append("meta.%s empty" % key)

    chapters = content.get("chapters", [])
    if len(chapters) != CHAPTER_COUNT:
        issues.append("chapter count %d != %d" % (len(chapters), CHAPTER_COUNT))
        return
    for i, chapter in enumerate(chapters):
        title = chapter.get("title", "")
        if title != CHAPTER_TITLES[i]:
            issues.append("chapter %d title %r != expected %r" % (i + 1, title, CHAPTER_TITLES[i]))
        if "研究与论文" in title:
            issues.append("forbidden chapter 研究与论文 found: %s" % title)
        if not str(chapter.get("statement", "")).strip():
            issues.append("chapter %d statement empty" % (i + 1))
        if not str(chapter.get("image", "")).strip():
            issues.append("chapter %d image missing" % (i + 1))
        items = chapter.get("items", [])
        if len(items) != ITEMS_PER_CHAPTER:
            issues.append("chapter %d has %d items, expected %d" % (i + 1, len(items), ITEMS_PER_CHAPTER))
        for j, item in enumerate(items):
            where = "chapter %d item %d" % (i + 1, j + 1)
            if len(item.get("title", "")) > MAX_TITLE_CHARS:
                issues.append("%s title too long (%d): %s" % (where, len(item["title"]), item["title"]))
            if len(item.get("summary", "")) > MAX_SUMMARY_CHARS:
                issues.append("%s summary too long (%d): %s" % (where, len(item["summary"]), item["title"]))
            kps = item.get("key_points", [])
            if not (3 <= len(kps) <= 4):
                issues.append("%s key_points count %d not in 3..4" % (where, len(kps)))
            for kp in kps:
                if len(kp) > MAX_KEY_POINT_CHARS:
                    issues.append("%s key point too long (%d): %s" % (where, len(kp), kp))
            if len(item.get("impact", "")) > MAX_IMPACT_CHARS:
                issues.append("%s impact too long (%d): %s" % (where, len(item["impact"]), item["title"]))
            for key in ("image", "link", "source", "meta"):
                if not str(item.get(key, "")).strip():
                    issues.append("%s.%s empty" % (where, key))


def check_cover(slide, content, slide_w, slide_h, issues, report_slides):
    meta = content["meta"]
    if shape_text(find(slide, "文本框 12")) != meta["title"]:
        issues.append("cover title mismatch")
    if shape_text(find(slide, "文本框 13")) != meta["publisher"]:
        issues.append("cover publisher mismatch")
    if not shape_text(find(slide, "文本框 1")):
        issues.append("cover date empty")
    table = find(slide, "表格 7")
    if table is None:
        issues.append("cover table missing")
    else:
        for (row, col) in [(0, 1), (0, 3), (1, 1), (1, 3), (2, 1), (2, 3)]:
            if not table.table.cell(row, col).text.strip():
                issues.append("cover table cell (%d,%d) empty" % (row, col))
    for sh in slide.shapes:
        check_bounds(sh, slide_w, slide_h, issues, "cover")
    report_slides.append({"index": 0, "role": "cover"})


def check_toc(slide, chapters, slide_w, slide_h, issues, report_slides):
    boxes = ["TextBox 2", "TextBox 3", "TextBox 4", "TextBox 6"]
    for box_no, chapter in zip(boxes, chapters):
        box = find(slide, box_no)
        if box is None or chapter["title"] not in shape_text(box):
            issues.append("global toc missing %s" % chapter["title"])
    for unused in ("TextBox 7", "TextBox 8"):
        box = find(slide, unused)
        if box is not None and shape_text(box):
            issues.append("unused global toc box %s should be empty" % unused)
    for sh in slide.shapes:
        check_bounds(sh, slide_w, slide_h, issues, "toc")
    report_slides.append({"index": 1, "role": "toc"})


def check_chapter(slide, chapter_no, chapter, slide_w, slide_h, issues, report_slides):
    index = len(report_slides)
    if shape_text(find(slide, "TextBox 4")) != chapter["title"]:
        issues.append("chapter %d title wrong" % chapter_no)
    if shape_text(find(slide, "TextBox 6")) != chapter["statement"]:
        issues.append("chapter %d statement wrong" % chapter_no)
    image = find(slide, "chapter_image")
    if image is None:
        issues.append("chapter %d image missing" % chapter_no)
        image_rect = None
    else:
        image_rect = rect_inches(image)
    if picture_count(slide) != 1:
        issues.append("chapter %d must have exactly 1 picture" % chapter_no)
    for idx in (1, 2):
        prefix = "chapter_preview_%d" % idx
        card = find(slide, prefix)
        if card is None:
            issues.append("chapter %d preview %d missing" % (chapter_no, idx))
            continue
        title_box = find(slide, prefix + "_title")
        summary_box = find(slide, prefix + "_summary")
        item = chapter["items"][idx - 1]
        if title_box is None or item["title"] not in shape_text(title_box):
            issues.append("chapter %d preview %d title wrong" % (chapter_no, idx))
        if summary_box is None or item["summary"] not in shape_text(summary_box):
            issues.append("chapter %d preview %d summary wrong" % (chapter_no, idx))
        for suffix in ("_badge", "_badge_text"):
            if find(slide, prefix + suffix) is None:
                issues.append("chapter %d preview %d missing %s" % (chapter_no, idx, suffix))
    for sh in slide.shapes:
        check_bounds(sh, slide_w, slide_h, issues, "chapter")
    report_slides.append({
        "index": index,
        "role": "chapter",
        "chapter": chapter_no,
        "image_rect": image_rect,
    })


def check_detail(slide, chapter_no, chapter, item_no, item, slide_w, slide_h,
                 issues, report_slides):
    index = len(report_slides)
    if shape_text(find(slide, "page_title")) != chapter["title"]:
        issues.append("detail page title wrong for %s" % item["title"])
    if not shape_text(find(slide, "page_kicker")):
        issues.append("detail page kicker empty for %s" % item["title"])
    if find(slide, "accent") is None:
        issues.append("detail accent missing for %s" % item["title"])
    if shape_text(find(slide, "detail_title")) != item["title"]:
        issues.append("detail title wrong for %s" % item["title"])
    if shape_text(find(slide, "detail_summary")) != item["summary"]:
        issues.append("detail summary wrong for %s" % item["title"])
    for kp_idx, kp in enumerate(item["key_points"], start=1):
        box = find(slide, "detail_kp_%d" % kp_idx)
        if box is None or kp not in shape_text(box):
            issues.append("detail key point %d wrong for %s" % (kp_idx, item["title"]))
    if shape_text(find(slide, "detail_impact")) != item["impact"]:
        issues.append("detail impact wrong for %s" % item["title"])
    if not shape_text(find(slide, "detail_source")):
        issues.append("detail source empty for %s" % item["title"])
    image = find(slide, "detail_image")
    if image is None:
        issues.append("detail image missing for %s" % item["title"])
        image_rect = None
    else:
        image_rect = rect_inches(image)
    if picture_count(slide) != 1:
        issues.append("detail page must have exactly 1 picture for %s" % item["title"])
    for sh in slide.shapes:
        check_bounds(sh, slide_w, slide_h, issues, "detail")
    report_slides.append({
        "index": index,
        "role": "detail",
        "chapter": chapter_no,
        "item": item_no,
        "image_rect": image_rect,
    })


def check_back(slide, slide_w, slide_h, issues, report_slides):
    if not shape_text(find(slide, "文本框 3")):
        issues.append("back cover slogan empty")
    for sh in slide.shapes:
        check_bounds(sh, slide_w, slide_h, issues, "back")
    report_slides.append({"index": len(report_slides), "role": "back"})


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--content", required=True)
    ap.add_argument("--report", required=True)
    args = ap.parse_args()

    with open(args.content, encoding="utf-8") as f:
        content = json.load(f)
    issues = []
    validate_content(content, issues)

    prs = Presentation(args.pptx)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    slides = list(prs.slides)
    if len(slides) != EXPECTED_SLIDES:
        issues.append("slide count %d != expected %d" % (len(slides), EXPECTED_SLIDES))

    report_slides = []
    if len(slides) > 0:
        check_cover(slides[0], content, slide_w, slide_h, issues, report_slides)
    if len(slides) > 1:
        check_toc(slides[1], content["chapters"], slide_w, slide_h, issues, report_slides)

    for chapter_no, chapter in enumerate(content.get("chapters", []), start=1):
        if len(slides) > len(report_slides):
            check_chapter(slides[len(report_slides)], chapter_no, chapter, slide_w,
                          slide_h, issues, report_slides)
        else:
            issues.append("missing chapter page %d" % chapter_no)
        for item_no, item in enumerate(chapter.get("items", []), start=1):
            if len(slides) > len(report_slides):
                check_detail(slides[len(report_slides)], chapter_no, chapter, item_no,
                             item, slide_w, slide_h, issues, report_slides)
            else:
                issues.append("missing detail page %s" % item.get("title", ""))

    if len(slides) > len(report_slides):
        check_back(slides[len(report_slides)], slide_w, slide_h, issues, report_slides)
    elif len(report_slides) < EXPECTED_SLIDES:
        issues.append("missing back cover")

    for slide in slides:
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                if "研究与论文" in sh.text_frame.text:
                    issues.append("forbidden text 研究与论文 on slide %d shape %r"
                                  % (slides.index(slide) + 1, sh.name))

    report = {
        "ok": not issues,
        "slide_count": len(slides),
        "expected_count": EXPECTED_SLIDES,
        "issues": issues,
        "slides": report_slides,
    }
    with open(args.report, "w", encoding="utf-8") as f:
        json.dump(report, f, ensure_ascii=False, indent=2)
    if issues:
        print("QA FAILED")
        for issue in issues:
            print(" -", issue)
        sys.exit(1)
    roles = ",".join(s["role"] for s in report_slides)
    print("QA OK: %d slides, roles=%s" % (len(slides), roles))


if __name__ == "__main__":
    main()
